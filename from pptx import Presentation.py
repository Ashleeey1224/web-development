from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Introduction and Overview
slide_1 = prs.slides.add_slide(prs.slide_layouts[1])
title_1 = slide_1.shapes.title
title_1.text = "Embracing Diversity in the Workplace"
subtitle_1 = slide_1.placeholders[1]
subtitle_1.text = "Understanding, Impact, and Strategies"

# Slide 2: Forms of Diversity
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
title_2.text = "Understanding Diversity"
content_2 = slide_2.placeholders[1].text_frame
content_2.text = "• Mental and Physical Ability: Definition, examples, and importance in the workplace\n• Culture and Race: Definition, examples, and importance in the workplace"

# Slide 3: Laws Dealing with Abuse or Discrimination
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
title_3.text = "Legal Frameworks for Protection"
content_3 = slide_3.placeholders[1].text_frame
content_3.text = "• Key anti-discrimination laws (e.g., Disability Discrimination Act, Racial Discrimination Act)\n• Protections against abuse and discrimination\n• Real-life examples of legal cases or precedents"

# Slide 4: Impact of Diversity on the Workplace
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
title_4.text = "Impact and Considerations"
content_4 = slide_4.placeholders[1].text_frame
content_4.text = "• Positive impacts of a diverse workforce (e.g., innovation, broader perspectives)\n• Special considerations for employees with diverse backgrounds (e.g., accommodations for disabilities, cultural sensitivity training)"

# Slide 5: Strategies and Comparison
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
title_5.text = "Strategies and International Comparison"
content_5 = slide_5.placeholders[1].text_frame
content_5.text = "• Strategies used to address diversity issues (e.g., diversity training, inclusive policies, behavioral changes)\n• Comparison of Australian practices of equal opportunity and social inclusion with those in your country of birth"

# Slide 6: Conclusion and Q&A
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
title_6.text = "Conclusion"
content_6 = slide_6.placeholders[1].text_frame
content_6.text = "• Summary of key points\n• Importance of continuous improvement in diversity and inclusion efforts\n• Invitation for questions and discussion"

# Save the presentation
pptx_file = "/mnt/data/Embracing_Diversity_in_the_Workplace.pptx"
prs.save(pptx_file)

pptx_file
